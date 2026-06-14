from __future__ import annotations

import ast
import html
import re
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Protocol

try:
    import fitz  # type: ignore
except ImportError:  # pragma: no cover
    fitz = None

try:
    from docx import Document as DocxDocument  # type: ignore
except ImportError:  # pragma: no cover
    DocxDocument = None

try:
    from openpyxl import load_workbook  # type: ignore
except ImportError:  # pragma: no cover
    load_workbook = None

try:
    from PIL import Image  # type: ignore
except ImportError:  # pragma: no cover
    Image = None

try:
    from pptx import Presentation  # type: ignore
except ImportError:  # pragma: no cover
    Presentation = None

try:
    import pytesseract  # type: ignore
except ImportError:  # pragma: no cover
    pytesseract = None

TEXT_CHUNK_SIZE = 4_000
TEXT_CHUNK_OVERLAP = 800
CODE_EXTENSIONS = {".js", ".ts", ".tsx", ".jsx", ".py", ".cs", ".java", ".go", ".rs"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp"}
OFFICE_EXTENSIONS = {".docx", ".pptx", ".xlsx"}
AUDIO_EXTENSIONS = {".mp3", ".mp4", ".wav", ".webm", ".ogg", ".m4a"}
VIDEO_EXTENSIONS = {".mp4", ".mov", ".mkv", ".webm"}


class IngestibleResource(Protocol):
    id: str
    name: str
    file_path: str


@dataclass
class ChunkPayload:
    content: str
    metadata: dict = field(default_factory=dict)


@dataclass
class ExtractedAsset:
    content: str
    metadata: dict = field(default_factory=dict)
    chunks: list[ChunkPayload] = field(default_factory=list)
    transcript: str | None = None


class ResourceProcessor(Protocol):
    def supports(self, resource: IngestibleResource, mime_type: str | None = None) -> bool: ...

    def extract(self, resource: IngestibleResource, mime_type: str | None = None) -> ExtractedAsset: ...


class AssetProcessingError(RuntimeError):
    pass


def chunk_text(text: str, *, base_metadata: dict | None = None) -> list[ChunkPayload]:
    normalized = " ".join(text.split())
    if not normalized:
        return []
    if len(normalized) <= TEXT_CHUNK_SIZE:
        return [ChunkPayload(content=normalized, metadata=base_metadata or {})]

    chunks: list[ChunkPayload] = []
    start = 0
    step = TEXT_CHUNK_SIZE - TEXT_CHUNK_OVERLAP
    while start < len(normalized):
        end = min(len(normalized), start + TEXT_CHUNK_SIZE)
        content = normalized[start:end].strip()
        if content:
            chunks.append(ChunkPayload(content=content, metadata=dict(base_metadata or {})))
        if end >= len(normalized):
            break
        start += step
    return chunks


class PdfProcessor:
    def supports(self, resource: IngestibleResource, mime_type: str | None = None) -> bool:
        return mime_type == "application/pdf" or resource.name.lower().endswith(".pdf")

    def extract(self, resource: IngestibleResource, mime_type: str | None = None) -> ExtractedAsset:
        path = Path(resource.file_path)
        if not path.exists() or fitz is None:
            return ExtractedAsset(content="")

        pages: list[str] = []
        chunks: list[ChunkPayload] = []
        page_count = 0
        with fitz.open(path) as document:
            page_count = document.page_count
            for page_index, page in enumerate(document, start=1):
                page_text = page.get_text("text").strip()
                if not page_text:
                    continue
                pages.append(page_text)
                chunks.extend(chunk_text(page_text, base_metadata={"page": page_index}))

        return ExtractedAsset(
            content="\n\n".join(pages).strip(),
            metadata={"page_count": page_count},
            chunks=chunks,
        )


class TextProcessor:
    supported_mime_types = {
        "text/plain",
        "text/markdown",
        "application/json",
        "application/xml",
    }

    def supports(self, resource: IngestibleResource, mime_type: str | None = None) -> bool:
        return (mime_type or "") in self.supported_mime_types or resource.name.lower().endswith((".txt", ".md", ".json", ".xml"))

    def extract(self, resource: IngestibleResource, mime_type: str | None = None) -> ExtractedAsset:
        path = Path(resource.file_path)
        text = path.read_text(encoding="utf-8", errors="ignore").strip() if path.exists() else ""
        return ExtractedAsset(content=text, chunks=chunk_text(text))


class CodeProcessor:
    def supports(self, resource: IngestibleResource, mime_type: str | None = None) -> bool:
        return Path(resource.name).suffix.lower() in CODE_EXTENSIONS

    def extract(self, resource: IngestibleResource, mime_type: str | None = None) -> ExtractedAsset:
        path = Path(resource.file_path)
        text = path.read_text(encoding="utf-8", errors="ignore") if path.exists() else ""
        extension = path.suffix.lower()
        chunks = self._chunk_python(text, resource.name) if extension == ".py" else self._chunk_generic(text, resource.name)
        return ExtractedAsset(
            content=text.strip(),
            metadata={"language": extension.lstrip(".")},
            chunks=chunks or chunk_text(text, base_metadata={"file_path": resource.name}),
        )

    def _chunk_python(self, text: str, filename: str) -> list[ChunkPayload]:
        try:
            tree = ast.parse(text)
        except SyntaxError:
            return self._chunk_generic(text, filename)

        lines = text.splitlines()
        chunks: list[ChunkPayload] = []
        for node in tree.body:
            if not isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef, ast.ClassDef)):
                continue
            start = getattr(node, "lineno", 1) - 1
            end = getattr(node, "end_lineno", start + 1)
            snippet = "\n".join(lines[start:end]).strip()
            if not snippet:
                continue
            chunks.append(
                ChunkPayload(
                    content=snippet,
                    metadata={
                        "file_path": filename,
                        "symbol": node.name,
                        "symbol_type": ("class" if isinstance(node, ast.ClassDef) else "function"),
                        "start_line": start + 1,
                        "end_line": end,
                    },
                )
            )
        return chunks

    def _chunk_generic(self, text: str, filename: str) -> list[ChunkPayload]:
        sections = re.split(r"\n(?=(?:export\s+)?(?:class|function|interface|type|struct|impl)\b)", text)
        return [
            ChunkPayload(
                content=section.strip(),
                metadata={"file_path": filename, "section_index": index},
            )
            for index, section in enumerate(sections)
            if section.strip()
        ]


class ImageProcessor:
    def supports(self, resource: IngestibleResource, mime_type: str | None = None) -> bool:
        return Path(resource.name).suffix.lower() in IMAGE_EXTENSIONS or (mime_type or "").startswith("image/")

    def extract(self, resource: IngestibleResource, mime_type: str | None = None) -> ExtractedAsset:
        path = Path(resource.file_path)
        extracted_text = ""
        if path.exists() and Image is not None and pytesseract is not None:
            try:
                extracted_text = pytesseract.image_to_string(Image.open(path)).strip()
            except Exception:  # noqa: BLE001
                extracted_text = ""
        description = f"Image asset named {resource.name}"
        combined = "\n\n".join(part for part in [description, extracted_text] if part).strip()
        return ExtractedAsset(
            content=combined,
            metadata={"ocr_text": extracted_text, "description": description},
            chunks=chunk_text(combined, base_metadata={"description": description}),
        )


class OfficeProcessor:
    def supports(self, resource: IngestibleResource, mime_type: str | None = None) -> bool:
        return Path(resource.name).suffix.lower() in OFFICE_EXTENSIONS

    def extract(self, resource: IngestibleResource, mime_type: str | None = None) -> ExtractedAsset:
        suffix = Path(resource.name).suffix.lower()
        if suffix == ".docx":
            return self._extract_docx(resource)
        if suffix == ".pptx":
            return self._extract_pptx(resource)
        if suffix == ".xlsx":
            return self._extract_xlsx(resource)
        return ExtractedAsset(content="")

    def _extract_docx(self, resource: IngestibleResource) -> ExtractedAsset:
        path = Path(resource.file_path)
        paragraphs: list[str] = []
        if DocxDocument is not None and path.exists():
            try:
                document = DocxDocument(path)
                paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
            except Exception:  # noqa: BLE001
                paragraphs = []
        elif path.exists():
            paragraphs = self._extract_docx_via_zip(path)
        text = "\n".join(paragraphs).strip()
        return ExtractedAsset(content=text, chunks=chunk_text(text))

    def _extract_docx_via_zip(self, path: Path) -> list[str]:
        try:
            with zipfile.ZipFile(path) as archive:
                document_xml = archive.read("word/document.xml").decode("utf-8", errors="ignore")
        except (FileNotFoundError, KeyError, zipfile.BadZipFile):
            return []
        paragraphs = re.findall(r"<w:t[^>]*>(.*?)</w:t>", document_xml)
        return [html.unescape(text).strip() for text in paragraphs if text.strip()]

    def _extract_pptx(self, resource: IngestibleResource) -> ExtractedAsset:
        path = Path(resource.file_path)
        if Presentation is None or not path.exists():
            return ExtractedAsset(content="")
        slides: list[str] = []
        chunks: list[ChunkPayload] = []
        try:
            presentation = Presentation(path)
            for slide_index, slide in enumerate(presentation.slides, start=1):
                texts = [getattr(shape, "text", "").strip() for shape in slide.shapes if getattr(shape, "text", "").strip()]
                slide_text = "\n".join(texts).strip()
                if not slide_text:
                    continue
                slides.append(slide_text)
                chunks.extend(chunk_text(slide_text, base_metadata={"slide": slide_index}))
        except Exception:  # noqa: BLE001
            return ExtractedAsset(content="")
        return ExtractedAsset(content="\n\n".join(slides).strip(), chunks=chunks)

    def _extract_xlsx(self, resource: IngestibleResource) -> ExtractedAsset:
        path = Path(resource.file_path)
        if load_workbook is None or not path.exists():
            return ExtractedAsset(content="")
        sheet_texts: list[str] = []
        chunks: list[ChunkPayload] = []
        try:
            workbook = load_workbook(path, read_only=True, data_only=True)
            for sheet in workbook.worksheets:
                rows: list[str] = []
                for row in sheet.iter_rows(values_only=True):
                    values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                    if values:
                        rows.append(" | ".join(values))
                if not rows:
                    continue
                sheet_text = "\n".join(rows)
                sheet_texts.append(f"{sheet.title}\n{sheet_text}")
                chunks.extend(chunk_text(sheet_text, base_metadata={"sheet": sheet.title}))
        except Exception:  # noqa: BLE001
            return ExtractedAsset(content="")
        return ExtractedAsset(content="\n\n".join(sheet_texts).strip(), chunks=chunks)


class AudioProcessor:
    def supports(self, resource: IngestibleResource, mime_type: str | None = None) -> bool:
        return (mime_type or "").startswith("audio/") or Path(resource.name).suffix.lower() in AUDIO_EXTENSIONS

    def extract(self, resource: IngestibleResource, mime_type: str | None = None) -> ExtractedAsset:
        transcript = _transcribe_media(resource)
        return ExtractedAsset(
            content=transcript,
            transcript=transcript,
            metadata={"transcript_available": bool(transcript)},
            chunks=chunk_text(transcript),
        )


class VideoProcessor:
    def supports(self, resource: IngestibleResource, mime_type: str | None = None) -> bool:
        return (mime_type or "").startswith("video/") or Path(resource.name).suffix.lower() in VIDEO_EXTENSIONS

    def extract(self, resource: IngestibleResource, mime_type: str | None = None) -> ExtractedAsset:
        transcript = _transcribe_media(resource)
        if not transcript:
            transcript = f"No spoken audio was detected in {resource.name}."
        return ExtractedAsset(
            content=transcript,
            transcript=transcript,
            metadata={"transcript_available": bool(transcript)},
            chunks=chunk_text(transcript),
        )


def _transcribe_media(resource: IngestibleResource) -> str:
    from app.services.whisper_transcription_service import WhisperTranscriptionService

    return WhisperTranscriptionService().transcribe_media(resource.file_path, resource.name)


class AssetProcessorRegistry:
    def __init__(self) -> None:
        self._processors: list[ResourceProcessor] = [
            PdfProcessor(),
            CodeProcessor(),
            AudioProcessor(),
            VideoProcessor(),
            OfficeProcessor(),
            ImageProcessor(),
            TextProcessor(),
        ]

    def select(self, resource: IngestibleResource, mime_type: str | None = None) -> ResourceProcessor:
        for processor in self._processors:
            if processor.supports(resource, mime_type):
                return processor
        return TextProcessor()


processor_registry = AssetProcessorRegistry()
